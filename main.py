import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import re

def apply_text_formatting(shape, text):
    """テキストにマークダウン装飾を適用"""
    text_frame = shape.text_frame
    text_frame.clear()  # 既存のテキストをクリア
    
    # 改行で分割して処理
    paragraphs = text.split('\n')
    
    for i, paragraph in enumerate(paragraphs):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        # 太字 (**text**) と斜体 (*text*) を処理
        parts = re.split(r'(\*\*.*?\*\*|\*.*?\*)', paragraph)
        
        for part in parts:
            run = p.add_run()
            
            # 太字処理
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
            # 斜体処理
            elif part.startswith('*') and part.endswith('*'):
                run.text = part[1:-1]
                run.font.italic = True
            # 通常テキスト
            else:
                run.text = part
        
        # フォント設定（日本語対応）
        for run in p.runs:
            run.font.name = 'Meiryo'
            run.font.size = Pt(14)

def markdown_to_pptx(markdown_text):
    prs = Presentation()
    
    # スライドを分割（---で区切る）
    slides_content = re.split(r'^---\s*$', markdown_text, flags=re.MULTILINE)
    
    for content in slides_content:
        if not content.strip():
            continue
            
        # タイトルスライドレイアウト
        slide_layout = prs.slide_layouts[1]  # タイトルとコンテンツ
        slide = prs.slides.add_slide(slide_layout)
        
        lines = content.split('\n')
        title = ""
        body_lines = []
        
        # タイトルと本文を分離
        for line in lines:
            if line.startswith('# '):
                title = line[2:].strip()
            else:
                if line.strip():
                    body_lines.append(line.strip())
        
        # タイトル設定
        if title:
            slide.shapes.title.text = title
        else:
            slide.shapes.title.text = "New Slide"
        
        # 本文設定
        if body_lines:
            content_shape = slide.placeholders[1]
            full_text = '\n'.join(body_lines)
            apply_text_formatting(content_shape, full_text)
            
            # 箇条書き自動処理
            text_frame = content_shape.text_frame
            for paragraph in text_frame.paragraphs:
                if paragraph.text.startswith('- '):
                    paragraph.level = 0
                    paragraph.text = paragraph.text[2:]
                elif paragraph.text.startswith('  - '):
                    paragraph.level = 1
                    paragraph.text = paragraph.text[4:]
    
    return prs

def save_pptx(prs):
    file_stream = BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return file_stream

# Streamlit UI
st.title('マークダウンからPowerPointへ変換（装飾対応版）')
st.markdown("""
以下のマークダウン記法が使用できます:
- `# タイトル` → スライドタイトル
- `---` → スライド区切り
- `**太字**` → 太字テキスト
- `*斜体*` → 斜体テキスト
- `- 項目` → 箇条書き
- `  - サブ項目` → ネストした箇条書き
""")

sample_md = """# プレゼンタイトル

**主要ポイント**
- *イノベーティブ*なソリューション
- **市場リーダー**としての地位
  - 業界シェアNo.1
  - 顧客満足度98%

---
## 第2スライド

*強調すべき*数字:
- 売上成長: **+30%**
- コスト削減: *25%*改善

---
### 詳細説明

このプロジェクトには以下の特徴があります:
- **高速処理**
  - 従来比*200%*の速度
- **使いやすいインターフェース**
"""

markdown_text = st.text_area(
    'マークダウンテキストを入力してください',
    height=300,
    value=sample_md
)

if st.button('PowerPointに変換'):
    if markdown_text:
        try:
            prs = markdown_to_pptx(markdown_text)
            pptx_file = save_pptx(prs)
            
            st.success('変換完了！')
            st.download_button(
                label='PowerPointをダウンロード',
                data=pptx_file,
                file_name='presentation.pptx',
                mime='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            
        except Exception as e:
            st.error(f'エラーが発生しました: {str(e)}')
    else:
        st.warning('マークダウンテキストを入力してください')